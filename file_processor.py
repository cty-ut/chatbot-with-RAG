import streamlit as st
import io
import docx
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import tempfile
import os
from PIL import Image
import base64
from openai import OpenAI
from config import GEMINI_API_KEY, GEMINI_BASE_URL, GEMINI_PICTURE_MODEL

# 创建Gemini客户端
gemini_client = OpenAI(
    api_key=GEMINI_API_KEY,
    base_url=GEMINI_BASE_URL
)


def process_image_file(file_obj, file_extension):
    """使用Gemini API处理图像文件"""
    temp_image_path = None
    extracted_text = ""

    try:
        # 保存上传的文件到临时文件
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as temp_img_file:
            temp_img_file.write(file_obj.read())
            temp_image_path = temp_img_file.name

        # 显示上传的图像
        image_obj = Image.open(temp_image_path)
        st.image(image_obj, caption="Uploaded Image", use_container_width=True)

        # 将图像编码为base64
        with open(temp_image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')

        # 构建多模态消息请求
        response = gemini_client.chat.completions.create(
            model=GEMINI_PICTURE_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Please analyze this image in detail. Describe what you see, including objects, people, settings, colors, and any notable details. Also, what does this image show or represent?"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{file_extension};base64,{base64_image}"
                            }
                        }
                    ]
                }
            ]
        )

        # 获取Gemini的图像分析结果
        if response.choices and response.choices[0].message:
            image_analysis = response.choices[0].message.content

            # 构建要发送给LLM的文本
            extracted_text = f"I've analyzed the uploaded image and here's what I found:\n\n{image_analysis}\n\nPlease provide more insights or answer any specific questions about this image."
        else:
            extracted_text = "I wasn't able to analyze the image properly. Please describe what you see in the image or ask specific questions about it."

    except Exception as e:
        st.error(f"Error processing image: {e}")
        extracted_text = f"There was an error processing the image: {str(e)}. Please try again or describe what you see in the image."
    finally:
        # 清理临时文件
        if temp_image_path and os.path.exists(temp_image_path):
            os.remove(temp_image_path)

    return extracted_text


def process_document_for_rag(file_obj):
    """处理上传的文档，提取文本用于RAG"""
    extracted_content = ""
    metadata = {"name": file_obj.name}
    file_extension = file_obj.name.split('.')[-1].lower()
    added_to_rag = False

    if file_extension == "pdf":
        current_doc_text = ""
        with io.BytesIO(file_obj.getvalue()) as pdf_stream:
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                current_doc_text += page.get_text() + "\n\n"
        extracted_content = current_doc_text

    elif file_extension == "docx":
        current_doc_text = ""
        word_document = docx.Document(file_obj)
        for para in word_document.paragraphs:
            current_doc_text += para.text + '\n'
        extracted_content = current_doc_text

    elif file_extension == "pptx":
        current_doc_text = ""
        presentation_doc = Presentation(file_obj)
        for slide_num, slide in enumerate(presentation_doc.slides, start=1):
            slide_text_content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text_content += shape.text + "\n"
            current_doc_text += slide_text_content + "\n----SLIDE BREAK----\n"
        extracted_content = current_doc_text

    elif file_extension == "txt":
        extracted_content = file_obj.getvalue().decode("utf-8")

    if extracted_content:
        st.session_state.rag_manager.add_document(extracted_content, metadata)
        added_to_rag = True

    if added_to_rag:
        return f"Processed {file_obj.name} and added to the knowledge base."
    else:
        return f"Could not extract content from {file_obj.name}, or the file is empty. Not added to the knowledge base."


def process_general_file(file_obj):
    """处理一般文件上传，返回提取的文本内容"""
    if not file_obj:
        return None

    file_extension = file_obj.name.split('.')[-1].lower()
    extracted_text = ""

    if file_extension == "docx":
        # 处理 Word 文件
        word_doc = docx.Document(file_obj)
        for para in word_doc.paragraphs:
            extracted_text += para.text + '\n'

    elif file_extension == "pdf":
        # 处理 PDF 文件
        with io.BytesIO(file_obj.getvalue()) as pdf_stream:
            pdf_doc = fitz.open(stream=pdf_stream, filetype="pdf")
            for page_num in range(pdf_doc.page_count):
                page_content = pdf_doc.load_page(page_num)
                extracted_text += f"Page {page_num + 1}:\n"
                extracted_text += page_content.get_text() + '\n\n'

    elif file_extension == "xlsx":
        # 处理 Excel 文件
        excel_dataframes = pd.read_excel(file_obj, sheet_name=None)
        for sheet_name, df_sheet in excel_dataframes.items():
            extracted_text += f"Sheet: {sheet_name}\n"
            extracted_text += df_sheet.to_string() + '\n\n'

    elif file_extension == "txt":
        # 处理文本文件
        extracted_text = file_obj.getvalue().decode("utf-8")

    elif file_extension == "pptx":
        # 处理 PPTX 文件
        ppt_doc = Presentation(file_obj)
        for slide_num, slide_item in enumerate(ppt_doc.slides, start=1):
            extracted_text += f"Slide {slide_num}:\n"
            for shape_item in slide_item.shapes:
                if hasattr(shape_item, "text"):
                    extracted_text += shape_item.text + "\n"
            extracted_text += "\n"


    elif file_extension in ["jpg", "jpeg", "png"]:

        extracted_text = process_image_file(file_obj, file_extension)

    elif file_extension in ["mp3", "wav", "m4a", "ogg"]:
        extracted_text = process_audio_file(file_obj, file_extension)

    return extracted_text


def process_audio_file(file_obj, file_extension):
    """处理音频文件，使用Gemini API进行转录和理解"""
    try:
        # 读取音频文件并进行base64编码
        audio_bytes = file_obj.getvalue()
        base64_audio = base64.b64encode(audio_bytes).decode('utf-8')

        # 显示音频播放器给用户
        st.audio(audio_bytes, format=f"audio/{file_extension}")

        with st.spinner("Transcribing and analyzing audio..."):
            # 调用Gemini API进行音频转录和理解
            response = gemini_client.chat.completions.create(
                model=GEMINI_PICTURE_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Please transcribe this audio and provide a summary of its content. If there are multiple speakers, try to distinguish them."
                            },
                            {
                                "type": "input_audio",
                                "input_audio": {
                                    "data": base64_audio,
                                    "format": file_extension
                                }
                            }
                        ]
                    }
                ]
            )

            # 提取Gemini的回复
            if response.choices and response.choices[0].message:
                audio_analysis = response.choices[0].message.content

                # 构建响应文本
                result_text = f"**Audio Analysis Results:**\n\n{audio_analysis}\n\nIs there anything specific about this audio content you'd like me to explain?"
                return result_text
            else:
                return "I couldn't analyze the audio properly. Is there something specific about the audio you'd like to ask about?"

    except Exception as e:
        st.error(f"Error processing audio: {e}")
        return f"There was an error processing the audio file: {str(e)}. Please try uploading again or describe the audio content."